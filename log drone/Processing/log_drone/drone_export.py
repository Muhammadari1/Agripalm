"""
Konversi JSON hasil download dari opendronelog_client.py (format "Open DroneLog Web
JSON Export": {_exportInfo, flight, telemetry, track, derived}) jadi rows telemetry
siap pakai, plus writer GPX & Excel.

Field JSON dikonfirmasi lewat eksplorasi manual (lihat catatan implementasi) terhadap
1 file log DJI Flip asli -- speed dalam telemetry.speed satuannya m/s (dicek: max(speed)
~= flight.maxSpeed, keduanya m/s), height = ketinggian relatif dari takeoff (dipakai
sebagai alt_m di rows, konsisten dengan drone_map.py), altitude = absolut/MSL (disimpan
terpisah kalau suatu saat dibutuhkan, tidak dipakai untuk peta).
"""
import json
import os
import shutil
import zipfile
from datetime import datetime, timedelta

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class DroneExportParseError(Exception):
    pass


def load_rows(json_path):
    """Baca file JSON hasil export opendronelog, kembalikan (rows, flight_meta).

    rows: list of dict per titik -- time_s, lat_deg, lon_deg, alt_m (height relatif),
    altitude_msl_m, speed_kmh, battery_pct, voltage, satellites. Baris dengan lat/lon
    (0,0) placeholder (sebelum GPS lock) dibuang.

    flight_meta: dict dari bagian "flight" JSON (aircraftName, droneSerial,
    batterySerial, startTime, durationSecs, totalDistance, dst.)
    """
    with open(json_path, encoding="utf-8") as f:
        data = json.load(f)

    flight = data.get("flight") or {}
    telemetry = data.get("telemetry") or {}

    for key in ("time", "latitude", "longitude"):
        if key not in telemetry:
            raise DroneExportParseError(f"Field telemetry '{key}' tidak ada di JSON hasil export.")

    n = len(telemetry["time"])

    def col(name):
        arr = telemetry.get(name)
        return arr if (arr is not None and len(arr) == n) else [None] * n

    times = col("time")
    lats = col("latitude")
    lons = col("longitude")
    heights = col("height")
    altitudes = col("altitude")
    speeds = col("speed")
    batteries = col("battery")
    voltages = col("batteryVoltage")
    satellites = col("satellites")

    rows = []
    for i in range(n):
        lat, lon = lats[i], lons[i]
        if lat is None or lon is None:
            continue
        if abs(lat) < 1e-6 and abs(lon) < 1e-6:
            continue  # placeholder sebelum GPS lock, bukan posisi asli
        rows.append({
            "sample_idx": i,
            "time_s": times[i],
            "lat_deg": lat,
            "lon_deg": lon,
            "alt_m": heights[i],
            "altitude_msl_m": altitudes[i],
            "speed_kmh": (speeds[i] * 3.6) if speeds[i] is not None else None,
            "battery_pct": batteries[i],
            "voltage": voltages[i],
            "satellites": satellites[i],
        })

    if not rows:
        raise DroneExportParseError("Tidak ada titik koordinat valid di JSON hasil export.")

    return rows, flight


# ---------------- GPX ----------------

def _parse_start_time(flight_meta):
    raw = flight_meta.get("startTime")
    if not raw:
        return None
    try:
        return datetime.fromisoformat(raw.replace("Z", "+00:00"))
    except ValueError:
        return None


def export_gpx(rows, flight_meta, out_path):
    """Tulis GPX 1.1 sederhana: 1 track, 1 segment, trkpt per baris (lat/lon/ele/time)."""
    start_time = _parse_start_time(flight_meta)
    name = flight_meta.get("displayName") or flight_meta.get("fileName") or "Flight"

    lines = [
        '<?xml version="1.0" encoding="UTF-8"?>',
        '<gpx version="1.1" creator="Agripalm Dashboard - Log Drone" '
        'xmlns="http://www.topografix.com/GPX/1/1">',
        "  <trk>",
        f"    <name>{_xml_escape(name)}</name>",
        "    <trkseg>",
    ]
    for r in rows:
        ele = r["alt_m"] if r["alt_m"] is not None else 0.0
        lines.append(f'      <trkpt lat="{r["lat_deg"]:.8f}" lon="{r["lon_deg"]:.8f}">')
        lines.append(f"        <ele>{ele:.2f}</ele>")
        if start_time is not None and r["time_s"] is not None:
            ts = start_time + timedelta(seconds=r["time_s"])
            lines.append(f'        <time>{ts.strftime("%Y-%m-%dT%H:%M:%SZ")}</time>')
        lines.append("      </trkpt>")
    lines.append("    </trkseg>")
    lines.append("  </trk>")
    lines.append("</gpx>")

    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def _xml_escape(text):
    return (
        text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        .replace('"', "&quot;").replace("'", "&apos;")
    )


# ---------------- Excel ----------------

TELEMETRY_COLUMNS = [
    "sample_idx", "time_s", "lat_deg", "lon_deg", "alt_m", "altitude_msl_m",
    "speed_kmh", "battery_pct", "voltage", "satellites",
]


def export_excel(rows, out_path):
    df = pd.DataFrame(rows, columns=TELEMETRY_COLUMNS)
    df.to_excel(out_path, index=False, engine="openpyxl")


# ---------------- Multi-file (Report Log Drone: gabungan banyak file jadi 1 rekap) ----

def export_gpx_multi(rows, file_metas, out_path):
    """Sama seperti export_gpx, tapi rows berasal dari banyak file sekaligus (key
    "track_idx" per row) -- tulis SATU blok <trk> per file (GPX 1.1 memang mendukung
    banyak track dalam 1 file secara native), masing-masing pakai flight_meta &
    nama file sendiri-sendiri supaya waktu <time> tiap titik tetap akurat per file.

    file_metas: dict {track_idx: (flight_meta, original_filename)}.
    """
    groups = {}
    order = []
    for r in rows:
        idx = r["track_idx"]
        if idx not in groups:
            groups[idx] = []
            order.append(idx)
        groups[idx].append(r)

    lines = [
        '<?xml version="1.0" encoding="UTF-8"?>',
        '<gpx version="1.1" creator="Agripalm Dashboard - Report Log Drone" '
        'xmlns="http://www.topografix.com/GPX/1/1">',
    ]
    for idx in order:
        pts = sorted(groups[idx], key=lambda p: (p["time_s"] is None, p["time_s"]))
        flight_meta, filename = file_metas.get(idx, ({}, f"Flight {idx + 1}"))
        start_time = _parse_start_time(flight_meta)
        name = flight_meta.get("displayName") or filename or f"Flight {idx + 1}"

        lines.append("  <trk>")
        lines.append(f"    <name>{_xml_escape(name)}</name>")
        lines.append("    <trkseg>")
        for r in pts:
            ele = r["alt_m"] if r["alt_m"] is not None else 0.0
            lines.append(f'      <trkpt lat="{r["lat_deg"]:.8f}" lon="{r["lon_deg"]:.8f}">')
            lines.append(f"        <ele>{ele:.2f}</ele>")
            if start_time is not None and r["time_s"] is not None:
                ts = start_time + timedelta(seconds=r["time_s"])
                lines.append(f'        <time>{ts.strftime("%Y-%m-%dT%H:%M:%SZ")}</time>')
            lines.append("      </trkpt>")
        lines.append("    </trkseg>")
        lines.append("  </trk>")
    lines.append("</gpx>")

    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


_MONTH_ABBR_EN = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]


def _ordinal_day(n):
    if 10 <= n % 100 <= 20:
        suffix = "th"
    else:
        suffix = {1: "st", 2: "nd", 3: "rd"}.get(n % 10, "th")
    return f"{n}{suffix}"


def _format_hms(total_seconds):
    """"HH:MM:SS" (jam SELALU 2 digit) -- dipakai untuk durasi per penerbangan."""
    total_seconds = int(round(total_seconds or 0))
    h, rem = divmod(total_seconds, 3600)
    m, s = divmod(rem, 60)
    return f"{h:02d}:{m:02d}:{s:02d}"


def _format_hms_total(total_seconds):
    """"H:MM:SS" (jam TIDAK di-zero-pad, bisa >24 jam) -- dipakai untuk baris total
    keseluruhan, beda gaya dari _format_hms supaya tidak salah baca sebagai jam
    hari (mis. total 26 jam ditulis "26:xx:xx", bukan wrap ke "02:xx:xx")."""
    total_seconds = int(round(total_seconds or 0))
    h, rem = divmod(total_seconds, 3600)
    m, s = divmod(rem, 60)
    return f"{h}:{m:02d}:{s:02d}"


def export_excel_monitoring(file_metas, pilot_name, drone_code, drone_sn, block_info, out_path):
    """Tabel rekap "MONITORING UTILISASI DRONE {kode}" gaya laporan operasional
    (GANTI dump telemetry mentah export_excel_multi yang lama) -- 1 baris per
    penerbangan/file dalam batch ini. Pakai openpyxl LANGSUNG (bukan
    pandas.to_excel) karena butuh merge cell manual: BULAN digabung per grup bulan
    yang sama, JUMLAH HARI/JUMLAH PENERBANGAN digabung per grup TANGGAL yang sama
    -- sesuai contoh laporan dari user.

    file_metas: dict {track_idx: (flight_meta, filename)}.
    block_info: dict {track_idx: {"estate":.., "blok":..}} dari
    drone_map.compute_flight_blocks() (boleh None/{} kalau Aresta tidak tersedia).
    pilot_name/drone_code: input manual user saat upload (form "Input"), BUKAN
    session username / serial number otomatis -- drone_sn cuma fallback kalau
    drone_code kosong."""
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter
    from Processing.log_drone.drone_map import _extract_flight_datetime

    block_info = block_info or {}

    flights = []
    for track_idx, (meta, filename) in file_metas.items():
        dt = _extract_flight_datetime(filename, meta.get("startTime"))
        if dt is None:
            continue
        dur = meta.get("durationSecs") or 0
        end_dt = dt + timedelta(seconds=dur)
        blk = block_info.get(track_idx) or {}
        flights.append({
            "date": dt.date(), "start": dt, "end": end_dt, "duration_s": dur,
            "region": blk.get("region") or "", "estate": blk.get("estate") or "",
            "blok": blk.get("blok") or "",
        })
    flights.sort(key=lambda f: f["start"])

    drone_label = drone_code or drone_sn or "-"
    headers = ["BULAN", "NAME", "KODE DRONE", "PIC", "REGION", "JUMLAH HARI",
               "JUMLAH PENERBANGAN", "ESTATE", "BLOK", "START_TIME", "END_TIME", "TOTAL_TIME"]
    n_cols = len(headers)

    wb = Workbook()
    ws = wb.active
    ws.title = "Monitoring"

    thin = Side(style="thin", color="B0B0B0")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    header_fill = PatternFill("solid", fgColor="ADD8E6")
    total_fill = PatternFill("solid", fgColor="D9D9D9")
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)

    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=n_cols)
    title_cell = ws.cell(row=1, column=1, value=f"MONITORING UTILISASI DRONE {drone_label}")
    title_cell.font = Font(bold=True, size=14)

    header_row = 3
    for c, h in enumerate(headers, start=1):
        cell = ws.cell(row=header_row, column=c, value=h)
        cell.font = Font(bold=True)
        cell.fill = header_fill
        cell.alignment = center
        cell.border = border

    widths = [8, 16, 14, 16, 10, 10, 12, 16, 20, 12, 12, 12]
    for c, w in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(c)].width = w

    if not flights:
        wb.save(out_path)
        return

    data_start = header_row + 1
    r = data_start
    prev_date = prev_month = None
    date_group_start = {}
    month_group_start = {}
    date_row_ranges = []
    month_row_ranges = []
    cell_center = Alignment(horizontal="center", vertical="center")

    for f in flights:
        d = f["date"]
        ym = (d.year, d.month)
        if d != prev_date:
            if prev_date is not None:
                date_row_ranges.append((prev_date, date_group_start[prev_date], r - 1))
            date_group_start[d] = r
            prev_date = d
        if ym != prev_month:
            if prev_month is not None:
                month_row_ranges.append((prev_month, month_group_start[prev_month], r - 1))
            month_group_start[ym] = r
            prev_month = ym

        ws.cell(row=r, column=2, value=f"{_MONTH_ABBR_EN[d.month - 1]}-{_ordinal_day(d.day)}-{d.year}")
        ws.cell(row=r, column=3, value=drone_label)
        ws.cell(row=r, column=4, value=pilot_name or "-")
        ws.cell(row=r, column=5, value=f["region"])
        ws.cell(row=r, column=8, value=f["estate"])
        ws.cell(row=r, column=9, value=f["blok"])
        ws.cell(row=r, column=10, value=f["start"].strftime("%H:%M:%S"))
        ws.cell(row=r, column=11, value=f["end"].strftime("%H:%M:%S"))
        ws.cell(row=r, column=12, value=_format_hms(f["duration_s"]))
        for c in range(1, n_cols + 1):
            ws.cell(row=r, column=c).border = border
            ws.cell(row=r, column=c).alignment = cell_center
        r += 1

    date_row_ranges.append((prev_date, date_group_start[prev_date], r - 1))
    month_row_ranges.append((prev_month, month_group_start[prev_month], r - 1))

    for d, r0, r1 in date_row_ranges:
        ws.cell(row=r0, column=6, value=1)
        ws.cell(row=r0, column=7, value=r1 - r0 + 1)
        if r1 > r0:
            ws.merge_cells(start_row=r0, start_column=6, end_row=r1, end_column=6)
            ws.merge_cells(start_row=r0, start_column=7, end_row=r1, end_column=7)

    for ym, r0, r1 in month_row_ranges:
        _, m = ym
        ws.cell(row=r0, column=1, value=_MONTH_ABBR_EN[m - 1])
        if r1 > r0:
            ws.merge_cells(start_row=r0, start_column=1, end_row=r1, end_column=1)

    total_row = r
    ws.merge_cells(start_row=total_row, start_column=1, end_row=total_row, end_column=n_cols - 1)
    total_label_cell = ws.cell(row=total_row, column=1, value="TOTAL FLIGHT TIME")
    total_seconds = sum(f["duration_s"] for f in flights)
    total_cell = ws.cell(row=total_row, column=n_cols, value=_format_hms_total(total_seconds))
    for cell in (total_label_cell, total_cell):
        cell.font = Font(bold=True)
        cell.fill = total_fill
        cell.alignment = cell_center
    for c in range(1, n_cols + 1):
        ws.cell(row=total_row, column=c).border = border

    wb.save(out_path)


def export_shp_multi(rows, file_metas, out_path):
    """Sama seperti export_gpx_multi, tapi format ESRI Shapefile: 1 LineString per
    file/track (EPSG:4326), atribut track_idx/filename/drone_sn/duration_s/
    distance_m. Shapefile bukan 1 file (.shp perlu .shx/.dbf/.prj menyertai),
    makanya di sini langsung DIZIP jadi 1 file supaya bisa didownload lewat
    browser sebagai 1 file utuh -- out_path harus berakhiran .zip."""
    groups = {}
    order = []
    for r in rows:
        idx = r["track_idx"]
        if idx not in groups:
            groups[idx] = []
            order.append(idx)
        groups[idx].append(r)

    import geopandas as gpd
    from shapely.geometry import LineString

    records = []
    for idx in order:
        pts = sorted(groups[idx], key=lambda p: (p["time_s"] is None, p["time_s"]))
        coords = [(p["lon_deg"], p["lat_deg"]) for p in pts]
        if len(coords) < 2:
            continue
        flight_meta, filename = file_metas.get(idx, ({}, f"Flight {idx + 1}"))
        records.append({
            "geometry": LineString(coords),
            "track_idx": idx,
            "filename": (filename or f"Flight {idx + 1}")[:80],
            "drone_sn": (flight_meta.get("droneSerial") or "")[:50],
            "duration_s": flight_meta.get("durationSecs"),
            "distance_m": flight_meta.get("totalDistance"),
        })

    if not records:
        raise DroneExportParseError("Tidak ada jalur valid untuk dibuat SHP.")

    gdf = gpd.GeoDataFrame(records, geometry="geometry", crs="EPSG:4326")

    tmp_dir = os.path.splitext(out_path)[0] + "_shp_tmp"
    if os.path.exists(tmp_dir):
        shutil.rmtree(tmp_dir, ignore_errors=True)
    os.makedirs(tmp_dir, exist_ok=True)
    try:
        shp_path = os.path.join(tmp_dir, "flight_tracks.shp")
        gdf.to_file(shp_path, driver="ESRI Shapefile")

        with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for fname in os.listdir(tmp_dir):
                zf.write(os.path.join(tmp_dir, fname), arcname=fname)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ---------------- PPTX Report (Report Log Drone: peta + rekap, 1 slide) --------------

def _format_minutes(duration_s):
    """Durasi dalam MENIT (bukan detik), 1 desimal, format koma Indonesia --
    dipakai di seluruh Report Log Drone (PPT & web) supaya konsisten."""
    if duration_s is None:
        return "-"
    return f"{duration_s / 60:.1f}".replace(".", ",") + " menit"


_C_NAVY = RGBColor(0x1E, 0x3A, 0x5F)
_C_DARK = RGBColor(0x1F, 0x29, 0x37)
_C_GRAY = RGBColor(0x6B, 0x72, 0x80)
_C_GREEN = RGBColor(0x16, 0xA3, 0x4A)
_C_CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)
_C_CARD_BORDER = RGBColor(0xDC, 0xDD, 0xE1)
_C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _format_distance_m(distance_m):
    if distance_m is None:
        return "-"
    return f"{distance_m:,.0f}".replace(",", ".") + " meter"


def _format_period_range(start_dt, end_dt):
    if start_dt is None:
        return "-"
    if end_dt is None or end_dt.date() == start_dt.date():
        end_str = end_dt.strftime("%H:%M") if end_dt else "-"
        return f"{start_dt.strftime('%d %b %Y')} ({start_dt.strftime('%H:%M')} - {end_str} WIB)"
    return f"{start_dt.strftime('%d %b %Y %H:%M')} - {end_dt.strftime('%d %b %Y %H:%M')} WIB"


def _add_box(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _set_text(shape, text, size, color, bold=False, align=None):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align is not None:
        p.alignment = align
    return p


def _add_section_header(slide, x, y, w, text):
    """Header seksi kecil: bar biru + judul bold -- dipakai buat "Peta Jalur
    Penerbangan & Area Blok", "Informasi Perangkat & Pilot", "Rincian Titik &
    Sesi Penerbangan"."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.02), Inches(0.045), Inches(0.22))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _C_NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    box = slide.shapes.add_textbox(Inches(x + 0.13), Inches(y), Inches(w - 0.13), Inches(0.26))
    _set_text(box, text, 13, _C_DARK, bold=True)


def export_pptx_report(out_path, batch_name, uploader_name, finished_at, upload_id,
                        duration_s, distance_m, drone_sn, drone_code, files_total,
                        flight_days_count, flight_start, flight_end,
                        estate_names, map_image_path, cluster_rows,
                        prs=None, region_name=None):
    """Slide dashboard "DRONE FLIGHT LOG REPORT" (gaya kartu statistik ala mockup
    user): header + badge nomor batch, 3 kartu statistik, panel kiri (peta ringkas
    + info perangkat/pilot), panel kanan (tabel rekap per area).

    uploader_name & drone_code SEKARANG input MANUAL user saat upload (form "Input"
    di drone_report_log.html -- kolom pilot_name/drone_code), BUKAN lagi session
    username / serial number JSON otomatis -- supaya cocok dipakai di laporan resmi
    (mis. "NEO SQA02", bukan serial number panjang). drone_sn (serial asli dari
    JSON) tetap disimpan, cuma dipakai fallback kalau drone_code kosong.

    flight_days_count: jumlah HARI UNIK terbang di batch ini (lihat
    drone_map.compute_flight_days_count), dipakai kartu "Total Penerbangan".

    cluster_rows: list of dict {titik, periode, uploader_name, duration_s,
    distance_m}, 1 per area/cluster (lihat drone_map.compute_cluster_recaps)."""
    owns_presentation = prs is None
    if owns_presentation:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    M = 0.35
    CONTENT_W = 13.333 - 2 * M

    # ── Header ──────────────────────────────────────────────────────────────
    report_title = "DRONE FLIGHT LOG REPORT"
    if region_name:
        report_title += f" - {region_name}"
    _set_text(slide.shapes.add_textbox(Inches(M), Inches(0.26), Inches(8.5), Inches(0.45)),
              report_title, 24, _C_NAVY, bold=True)
    _set_text(slide.shapes.add_textbox(Inches(M), Inches(0.68), Inches(8.5), Inches(0.3)),
              "Laporan Hasil Pemetaan & Operasional Penerbangan Drone", 11, _C_GRAY)

    batch_code = f"#{finished_at:%Y-%m%d}-{upload_id}" if finished_at else f"#{upload_id}"
    badge = _add_box(slide, 10.35, 0.28, 2.633, 0.34, _C_NAVY)
    _set_text(badge, f"BATCH LOG {batch_code}", 11, _C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    uploaded_str = finished_at.strftime("%d %b %Y, %H:%M") if finished_at else "-"
    _set_text(slide.shapes.add_textbox(Inches(10.35), Inches(0.64), Inches(2.633), Inches(0.24)),
              f"Tanggal Unggah: {uploaded_str} WIB", 9, _C_GRAY, align=PP_ALIGN.RIGHT)

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(M), Inches(1.0), Inches(CONTENT_W), Inches(0.018))
    divider.fill.solid(); divider.fill.fore_color.rgb = _C_CARD_BORDER
    divider.line.fill.background(); divider.shadow.inherit = False

    # ── Kartu statistik ─────────────────────────────────────────────────────
    card_gap = 0.15
    card_w = (CONTENT_W - 2 * card_gap) / 3
    card_y = 1.15
    card_h = 0.85

    session_sub = f"{files_total} Sesi Flight" if files_total else ""
    cards = [
        ("TOTAL DURASI TERBANG", _format_minutes(duration_s), session_sub),
        ("TOTAL JARAK TEMPUH", _format_distance_m(distance_m), ""),
        ("TOTAL PENERBANGAN", f"{flight_days_count} Hari" if flight_days_count else "-", ""),
    ]
    for i, (caption, value, sub) in enumerate(cards):
        x = M + i * (card_w + card_gap)
        _add_box(slide, x, card_y, card_w, card_h, _C_CARD_BG, _C_CARD_BORDER)
        _set_text(slide.shapes.add_textbox(Inches(x + 0.12), Inches(card_y + 0.08), Inches(card_w - 0.24), Inches(0.22)),
                  caption, 8.5, _C_GRAY, bold=True)
        _set_text(slide.shapes.add_textbox(Inches(x + 0.12), Inches(card_y + 0.28), Inches(card_w - 0.24), Inches(0.34)),
                  value, 19, _C_DARK, bold=True)
        if sub:
            _set_text(slide.shapes.add_textbox(Inches(x + 0.12), Inches(card_y + 0.62), Inches(card_w - 0.24), Inches(0.2)),
                      sub, 8.5, _C_GREEN)

    # ── Panel kiri: peta + info perangkat/pilot ─────────────────────────────
    left_x = M
    left_w = 6.15
    y = 2.15
    _add_section_header(slide, left_x, y, left_w, "Peta Jalur Penerbangan & Area Blok")
    y += 0.32
    map_top = y
    if map_image_path and os.path.exists(map_image_path):
        pic = slide.shapes.add_picture(map_image_path, Inches(left_x), Inches(map_top), width=Inches(left_w))
        y = map_top + pic.height / 914400 + 0.14  # EMU -> inch
    else:
        y += 0.2

    _add_section_header(slide, left_x, y, left_w, "Informasi Perangkat & Pilot")
    y += 0.32
    period_str = _format_period_range(flight_start, flight_end)
    lokasi_str = "Blok " + " / ".join(estate_names) if estate_names else "-"
    info_lines = [
        ("Nama Pilot/User", uploader_name or "-"),
        ("Kode Drone", drone_code or drone_sn or "-"),
        ("Waktu Mulai - Selesai", period_str),
        ("Lokasi/Area Kerja", lokasi_str),
    ]
    for label, value in info_lines:
        box = slide.shapes.add_textbox(Inches(left_x), Inches(y), Inches(left_w), Inches(0.24))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{label}"
        r1.font.size = Pt(10); r1.font.color.rgb = _C_GRAY
        r2 = p.add_run(); r2.text = f" : {value}"
        r2.font.size = Pt(10); r2.font.color.rgb = _C_DARK; r2.font.bold = True
        y += 0.27

    # ── Panel kanan: tabel rekap per area ────────────────────────────────────
    right_x = left_x + left_w + 0.35
    right_w = 13.333 - M - right_x
    y = 2.15
    _add_section_header(slide, right_x, y, right_w, "Rincian Titik & Sesi Penerbangan")
    y += 0.32
    table_top = y

    headers = ["No", "Kode Titik", "Operator", "Periode Waktu", "Durasi", "Jarak"]
    n_rows = len(cluster_rows) + 1
    table_max_h = 7.15 - table_top
    table_shape = slide.shapes.add_table(n_rows, len(headers), Inches(right_x), Inches(table_top),
                                          Inches(right_w), Inches(min(table_max_h, 0.35 * n_rows)))
    table = table_shape.table
    table.first_row = False  # kita atur warna header manual, jangan pakai style bawaan
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = _C_NAVY
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = _C_WHITE
        if c >= 4:
            p.alignment = PP_ALIGN.RIGHT

    for i, cr in enumerate(cluster_rows, start=1):
        dist = cr.get("distance_m")
        values = [
            str(i), str(cr.get("titik") or "-"), str(cr.get("uploader_name") or "-"),
            str(cr.get("periode") or "-"),
            _format_minutes(cr.get("duration_s")),
            f"{dist:,.0f}".replace(",", ".") + " m" if dist is not None else "-",
        ]
        for c, val in enumerate(values):
            cell = table.cell(i, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _C_WHITE if i % 2 else _C_CARD_BG
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10.5); p.font.color.rgb = _C_DARK
            if c >= 4:
                p.alignment = PP_ALIGN.RIGHT

    if owns_presentation:
        prs.save(out_path)
